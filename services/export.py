"""
Document export helpers — PDF, DOCX, PPTX generation.

All functions are pure (no app/DB dependencies) and can be called from
any router or service.
"""

import html
import io
import re
from typing import Any

from services.text_svc import _normalize_hex_color, _hex_to_rgb


# ---------------------------------------------------------------------------
# Markdown table splitting (shared by PDF/DOCX builders)
# ---------------------------------------------------------------------------


def _split_markdown_tables(text: str) -> list[tuple[str, str]]:
    """Split markdown text into ('text', content) and ('table', rows) chunks."""
    lines = text.split("\n")
    chunks: list[tuple[str, str]] = []
    current_text: list[str] = []
    current_table: list[str] = []

    def flush_text():
        if current_text:
            chunks.append(("text", "\n".join(current_text)))
            current_text.clear()

    def flush_table():
        if current_table:
            chunks.append(("table", "\n".join(current_table)))
            current_table.clear()

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            if not current_table:
                flush_text()
            current_table.append(stripped)
        else:
            if current_table:
                flush_table()
            current_text.append(line)
    flush_text()
    flush_table()
    return chunks


# ---------------------------------------------------------------------------
# ReportLab helpers
# ---------------------------------------------------------------------------


def _build_reportlab_table(table_md, page_width, cell_style, header_style,
                            primary_color, secondary_color, colors_mod,
                            TableCls, TableStyleCls, ParagraphCls):
    """Convert markdown table lines into a ReportLab Table with text wrapping."""
    rows_raw = [line.strip("|").split("|") for line in table_md.split("\n") if line.strip()]
    rows = []
    for cells in rows_raw:
        cleaned = [c.strip() for c in cells]
        if all(set(c) <= set("- :") for c in cleaned if c):
            continue
        rows.append(cleaned)
    if not rows:
        return None
    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")
    col_width = page_width / max_cols
    table_data = []
    for i, row in enumerate(rows):
        style = header_style if i == 0 else cell_style
        table_data.append([ParagraphCls(html.escape(cell), style) for cell in row])
    t = TableCls(table_data, colWidths=[col_width] * max_cols, repeatRows=1)
    style_commands = [
        ("GRID", (0, 0), (-1, -1), 0.5, colors_mod.Color(0.7, 0.7, 0.7)),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ("BACKGROUND", (0, 0), (-1, 0), primary_color),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors_mod.white),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors_mod.white, colors_mod.Color(0.97, 0.97, 0.96)]),
    ]
    t.setStyle(TableStyleCls(style_commands))
    return t


def _markdown_to_reportlab(text: str) -> str:
    """Convert common markdown patterns to reportlab-compatible XML."""
    lines = text.split("\n")
    result: list[str] = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            result.append("<br/>")
            continue
        if re.match(r"^#{1,4}\s+", stripped):
            stripped = re.sub(r"^#{1,4}\s+", "", stripped)
            stripped = f"<b>{html.escape(stripped)}</b>"
            result.append(stripped)
            continue
        if stripped.startswith("- ") or stripped.startswith("* "):
            bullet_text = stripped[2:].strip()
            bullet_text = html.escape(bullet_text)
            bullet_text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", bullet_text)
            bullet_text = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"<i>\1</i>", bullet_text)
            result.append(f"&nbsp;&nbsp;&bull;&nbsp;{bullet_text}")
            continue
        if re.match(r"^\d+\.\s+", stripped):
            num_match = re.match(r"^(\d+\.)\s+(.*)", stripped)
            if num_match:
                num = html.escape(num_match.group(1))
                rest = html.escape(num_match.group(2))
                rest = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", rest)
                rest = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"<i>\1</i>", rest)
                result.append(f"&nbsp;&nbsp;{num}&nbsp;{rest}")
                continue
        if stripped.startswith("|") and stripped.endswith("|"):
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            if all(set(c) <= set("- :") for c in cells):
                continue
            row = "&nbsp;&nbsp;|&nbsp;&nbsp;".join(html.escape(c) for c in cells)
            result.append(f"<font face='Courier'>{row}</font>")
            continue
        escaped = html.escape(stripped)
        escaped = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", escaped)
        escaped = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"<i>\1</i>", escaped)
        result.append(escaped)
    return "<br/>".join(result)


def build_pdf_bytes(
    title: str,
    sections: list[dict[str, str]],
    branding: dict[str, Any] | None = None,
    logo_bytes: bytes | None = None,
) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    from reportlab.graphics.shapes import Drawing, Line

    styles = getSampleStyleSheet()
    branding = branding or {}
    primary_hex = _normalize_hex_color(branding.get("primary_color")) or "#1a1714"
    secondary_hex = _normalize_hex_color(branding.get("secondary_color")) or "#8c5a24"
    primary_color = colors.HexColor(primary_hex)
    secondary_color = colors.HexColor(secondary_hex)
    title_style = ParagraphStyle("BrandTitle", parent=styles["Title"], textColor=primary_color, fontSize=20, leading=24)
    heading_style = ParagraphStyle("BrandHeading", parent=styles["Heading2"], textColor=secondary_color, fontSize=14, leading=18, spaceBefore=14)
    subtitle_style = ParagraphStyle("BrandSubtitle", parent=styles["BodyText"], textColor=secondary_color, fontSize=10)
    body_style = ParagraphStyle("BrandBody", parent=styles["BodyText"], fontSize=10, leading=14, spaceBefore=2, spaceAfter=4)
    buffer = io.BytesIO()
    page_width = LETTER[0] - 108
    doc = SimpleDocTemplate(buffer, pagesize=LETTER, topMargin=54, bottomMargin=54, leftMargin=54, rightMargin=54)
    story: list[Any] = []

    title_para = Paragraph(html.escape(title), title_style)
    brand_label = branding.get("brand_name") or ""
    if logo_bytes:
        try:
            logo_img = Image(io.BytesIO(logo_bytes), width=72, height=72, kind="proportional")
            header_data = [[logo_img, title_para]]
            header_table = Table(header_data, colWidths=[90, page_width - 90])
            header_table.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
                ("TOPPADDING", (0, 0), (-1, -1), 0),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
            ]))
            story.append(header_table)
        except Exception:
            story.append(title_para)
    else:
        story.append(title_para)

    if brand_label:
        story.append(Spacer(1, 4))
        story.append(Paragraph(html.escape(brand_label), subtitle_style))

    story.append(Spacer(1, 10))
    d = Drawing(page_width, 2)
    d.add(Line(0, 1, page_width, 1, strokeColor=primary_color, strokeWidth=1.5))
    story.append(d)
    story.append(Spacer(1, 14))

    table_cell_style = ParagraphStyle("TableCell", parent=body_style, fontSize=8, leading=10, spaceBefore=0, spaceAfter=0)
    table_header_style = ParagraphStyle("TableHeader", parent=table_cell_style, textColor=colors.white)

    for section in sections:
        heading = (section.get("heading") or "").strip()
        body = (section.get("body") or "").strip()
        if heading:
            story.append(Paragraph(html.escape(heading), heading_style))
            story.append(Spacer(1, 6))
        if body:
            chunks = _split_markdown_tables(body)
            for chunk_type, chunk_content in chunks:
                if chunk_type == "table":
                    table_obj = _build_reportlab_table(
                        chunk_content, page_width, table_cell_style, table_header_style,
                        primary_color, secondary_color, colors,
                        Table, TableStyle, Paragraph,
                    )
                    if table_obj:
                        story.append(table_obj)
                        story.append(Spacer(1, 8))
                else:
                    converted = _markdown_to_reportlab(chunk_content)
                    if converted.strip():
                        story.append(Paragraph(converted, body_style))
                        story.append(Spacer(1, 8))

    doc.build(story)
    return buffer.getvalue()


def build_pptx_bytes(
    title: str,
    slides: list[dict[str, Any]],
    branding: dict[str, Any] | None = None,
    logo_bytes: bytes | None = None,
) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    branding = branding or {}
    primary_rgb = RGBColor(*_hex_to_rgb(branding.get("primary_color")))
    secondary_rgb = RGBColor(*_hex_to_rgb(branding.get("secondary_color")))
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_fill = title_slide.background.fill
    title_fill.solid()
    title_fill.fore_color.rgb = primary_rgb
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle = title_slide.placeholders[1]
    subtitle.text = f'Branded for {branding.get("brand_name") or "Meeting Analyzer"}'
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(245, 240, 235)
    if logo_bytes:
        try:
            title_slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(9.0), Inches(0.3), height=Inches(0.7))
        except Exception:
            pass
    for slide_data in slides[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = (slide_data.get("title") or "Slide").strip()
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_rgb
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = slide_data.get("bullets") or []
        for idx, bullet in enumerate(bullets[:6]):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = str(bullet).strip()
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = secondary_rgb
        if logo_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(9.1), Inches(0.2), height=Inches(0.45))
            except Exception:
                pass
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _add_markdown_runs(paragraph, text: str, _re) -> None:
    """Add runs to a docx paragraph, converting **bold** and *italic* markers."""
    from docx.shared import Pt
    parts = _re.split(r"(\*\*[^*]+?\*\*|\*[^*]+?\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
            run.font.size = Pt(10)
        elif part.startswith("*") and part.endswith("*") and not part.startswith("**"):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
            run.font.size = Pt(10)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(10)


def _add_docx_table(doc, table_md, primary_rgb, secondary_rgb, _re, Pt, RGBColor, OxmlElement, qn):
    """Convert markdown table lines into a python-docx table with text wrapping."""
    rows_raw = [line.strip("|").split("|") for line in table_md.split("\n") if line.strip()]
    rows = []
    for cells in rows_raw:
        cleaned = [c.strip() for c in cells]
        if all(set(c) <= set("- :") for c in cleaned if c):
            continue
        rows.append(cleaned)
    if not rows:
        return
    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")
    table = doc.add_table(rows=len(rows), cols=max_cols)
    table.style = "Table Grid"
    table.autofit = True
    for i, row_data in enumerate(rows):
        row = table.rows[i]
        for j, cell_text in enumerate(row_data):
            cell = row.cells[j]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(cell_text)
            run.font.size = Pt(9)
            if i == 0:
                run.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                shading = OxmlElement("w:shd")
                shading.set(qn("w:val"), "clear")
                shading.set(qn("w:color"), "auto")
                shading.set(qn("w:fill"), str(primary_rgb).replace("#", ""))
                cell._tc.get_or_add_tcPr().append(shading)
    doc.add_paragraph("")


def build_docx_bytes(
    title: str,
    sections: list[dict[str, str]],
    branding: dict[str, Any] | None = None,
    logo_bytes: bytes | None = None,
) -> bytes:
    import re as _re
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    branding = branding or {}
    primary_rgb = RGBColor(*_hex_to_rgb(branding.get("primary_color")))
    secondary_rgb = RGBColor(*_hex_to_rgb(branding.get("secondary_color")))

    if logo_bytes:
        try:
            doc.add_picture(io.BytesIO(logo_bytes), height=Inches(0.8))
        except Exception:
            pass

    title_para = doc.add_heading(title, level=0)
    for run in title_para.runs:
        run.font.color.rgb = primary_rgb

    brand_name = branding.get("brand_name")
    if brand_name:
        subtitle_para = doc.add_paragraph(brand_name)
        subtitle_para.style = doc.styles["Subtitle"]
        for run in subtitle_para.runs:
            run.font.color.rgb = secondary_rgb

    rule_para = doc.add_paragraph()
    pPr = rule_para._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), (branding.get("primary_color") or "#1a1714").lstrip("#"))
    pBdr.append(bottom)
    pPr.append(pBdr)

    for section in sections:
        heading = (section.get("heading") or "").strip()
        body = (section.get("body") or "").strip()
        if heading:
            h_para = doc.add_heading(heading, level=1)
            for run in h_para.runs:
                run.font.color.rgb = secondary_rgb
        if body:
            chunks = _split_markdown_tables(body)
            for chunk_type, chunk_content in chunks:
                if chunk_type == "table":
                    _add_docx_table(doc, chunk_content, primary_rgb, secondary_rgb, _re, Pt, RGBColor, OxmlElement, qn)
                else:
                    for line in chunk_content.split("\n"):
                        stripped = line.strip()
                        if not stripped:
                            doc.add_paragraph("")
                            continue
                        if stripped.startswith("- ") or stripped.startswith("* "):
                            bullet_text = stripped[2:].strip()
                            p = doc.add_paragraph(style="List Bullet")
                            _add_markdown_runs(p, bullet_text, _re)
                        elif _re.match(r"^\d+\.\s+", stripped):
                            num_match = _re.match(r"^(\d+\.)\s+(.*)", stripped)
                            if num_match:
                                p = doc.add_paragraph(style="List Number")
                                _add_markdown_runs(p, num_match.group(2), _re)
                        elif _re.match(r"^#{1,4}\s+", stripped):
                            sub_heading = _re.sub(r"^#{1,4}\s+", "", stripped)
                            h = doc.add_heading(sub_heading, level=2)
                            for run in h.runs:
                                run.font.color.rgb = secondary_rgb
                        else:
                            p = doc.add_paragraph()
                            _add_markdown_runs(p, stripped, _re)

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()
