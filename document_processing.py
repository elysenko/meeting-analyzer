"""Document text extraction, table extraction, and MIME detection helpers."""
import io
import logging
import mimetypes
import os

logger = logging.getLogger("meeting-analyzer")

def _detect_mime(filename: str) -> str:
    mime, _ = mimetypes.guess_type(filename)
    return mime or "application/octet-stream"


def _extract_tables_pdf(data: bytes) -> list[dict]:
    import pdfplumber
    tables = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            for j, table in enumerate(page.extract_tables()):
                if not table or len(table) < 2:
                    continue
                headers = [str(c or "") for c in table[0]]
                rows = [[str(c or "") for c in row] for row in table[1:]]
                tables.append({"name": f"Page {i+1}, Table {j+1}", "page": i + 1, "headers": headers, "rows": rows})
    return tables


def _extract_tables_docx(data: bytes) -> list[dict]:
    from docx import Document
    doc = Document(io.BytesIO(data))
    tables = []
    for i, tbl in enumerate(doc.tables):
        rows_data = []
        for row in tbl.rows:
            rows_data.append([cell.text.strip() for cell in row.cells])
        if len(rows_data) < 2:
            continue
        tables.append({"name": f"Table {i+1}", "page": None, "headers": rows_data[0], "rows": rows_data[1:]})
    return tables


def _extract_tables_xlsx(data: bytes) -> list[dict]:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    tables = []
    for ws in wb.worksheets:
        all_rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):
                all_rows.append(cells)
        if len(all_rows) < 2:
            continue
        tables.append({"name": ws.title, "page": None, "headers": all_rows[0], "rows": all_rows[1:]})
    return tables


def _extract_tables_csv(data: bytes) -> list[dict]:
    import csv
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    all_rows = [row for row in reader if any(c.strip() for c in row)]
    if len(all_rows) < 2:
        return []
    return [{"name": "CSV", "page": None, "headers": all_rows[0], "rows": all_rows[1:]}]


def _extract_tables_sync(data: bytes, mime_type: str, filename: str) -> list[dict]:
    ext = os.path.splitext(filename.lower())[1]
    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            return _extract_tables_pdf(data)
        elif ext == ".docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_tables_docx(data)
        elif ext in (".xlsx", ".xls") or mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            return _extract_tables_xlsx(data)
        elif ext == ".csv" or mime_type == "text/csv":
            return _extract_tables_csv(data)
        else:
            return []
    except Exception as e:
        logger.warning("Table extraction failed for %s: %s", filename, e)
        return []


def _ocr_pdf_pages(data: bytes) -> str | None:
    """Render PDF pages to images and OCR each one with Tesseract."""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
        dpi = 150 if len(data) > 2_000_000 else 200
        images = convert_from_bytes(data, dpi=dpi)
        parts = []
        for img in images:
            text = pytesseract.image_to_string(img).strip()
            if text:
                parts.append(text)
            del img
        return "\n".join(parts) if parts else None
    except Exception as e:
        logger.warning("PDF OCR failed: %s", e)
        return None


def _extract_text_sync(data: bytes, mime_type: str, filename: str) -> str | None:
    """Extract text from document bytes. Returns None if unsupported."""
    ext = os.path.splitext(filename.lower())[1]
    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
            result = "\n".join(parts) if parts else ""
            if len(result.strip()) < 50:
                ocr_text = _ocr_pdf_pages(data)
                if ocr_text and len(ocr_text.strip()) > len(result.strip()):
                    return ocr_text
            return result
        elif ext == ".docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        elif ext in (".xlsx", ".xls") or mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            return "\n".join(rows)
        elif ext == ".pptx" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Word-only fallback — for OCR use _extract_text_pptx_hybrid_sync instead
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                parts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                parts.append("\t".join(cells))
            return "\n".join(parts)
        elif ext == ".doc" or mime_type == "application/msword":
            return _extract_text_doc_sync(data, filename)
        elif ext in (".txt", ".md", ".csv", ".log", ".json", ".yaml", ".yml") or (
            mime_type and mime_type.startswith("text/")
        ):
            return data.decode("utf-8", errors="replace")
        else:
            return None
    except Exception as e:
        logger.warning("Text extraction failed for %s: %s", filename, e)
        return None


def _extract_text_doc_sync(data: bytes, filename: str) -> str | None:
    """Extract text from legacy .doc files via LibreOffice → PDF → pypdf."""
    import subprocess, tempfile
    ext = os.path.splitext(filename.lower())[1] or ".doc"
    try:
        with tempfile.TemporaryDirectory(prefix="lo-doc-") as tmpdir:
            src = os.path.join(tmpdir, f"input{ext}")
            with open(src, "wb") as f:
                f.write(data)
            env = os.environ.copy()
            env["HOME"] = tmpdir
            result = subprocess.run(
                ["soffice", "--headless", "--norestore", "--nofirststartwizard",
                 "--convert-to", "pdf", "--outdir", tmpdir, src],
                capture_output=True, timeout=120, env=env,
            )
            if result.returncode != 0:
                logger.warning("soffice .doc→pdf failed for %s: %s", filename, result.stderr.decode()[-300:])
                return None
            pdf_path = os.path.join(tmpdir, "input.pdf")
            if not os.path.exists(pdf_path):
                logger.warning("soffice ran but no PDF output for %s", filename)
                return None
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(pdf_bytes))
        parts = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n".join(parts) if parts else None
    except Exception as e:
        logger.warning("DOC text extraction failed for %s: %s", filename, e)
        return None


def _extract_text_pptx_hybrid_sync(data: bytes, filename: str, pdf_bytes: bytes | None = None) -> str | None:
    """Hybrid PPTX extraction: python-pptx for text frames + Tesseract OCR on rendered slides.

    pdf_bytes should be the already-converted LibreOffice PDF to avoid running
    the conversion twice. When provided, each slide page is rendered at 200 DPI
    and OCR'd; the richer result (per slide) is used.
    """
    # ── Step 1: python-pptx direct text extraction ─────────────────────────
    pptx_slides: list[str] = []
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append("\t".join(cells))
            pptx_slides.append("\n".join(parts))
    except Exception as e:
        logger.warning("python-pptx extraction failed for %s: %s", filename, e)

    # ── Step 2: Tesseract OCR on rendered PDF pages ─────────────────────────
    ocr_slides: list[str] = []
    if pdf_bytes:
        try:
            from pdf2image import convert_from_bytes
            import pytesseract
            # 150 DPI: ~1.5 MB per slide vs ~4 MB at 200 DPI — reduces peak memory
            images = convert_from_bytes(pdf_bytes, dpi=150)
            for img in images:
                if img.mode in ("RGBA", "LA"):
                    img = img.convert("RGB")
                text = pytesseract.image_to_string(img, config="--psm 3").strip()
                ocr_slides.append(text)
                del img
        except Exception as e:
            logger.warning("PPTX slide OCR failed for %s: %s", filename, e)

    # ── Merge: per-slide winner, or whichever source is richer overall ──────
    if not pptx_slides and not ocr_slides:
        return None
    if pptx_slides and ocr_slides and len(pptx_slides) == len(ocr_slides):
        merged: list[str] = []
        for ps, os_ in zip(pptx_slides, ocr_slides):
            merged.append(os_ if len(os_) > len(ps) * 1.2 else (ps or os_))
        return "\n\n".join(s for s in merged if s) or None
    pptx_text = "\n\n".join(s for s in pptx_slides if s)
    ocr_text = "\n\n".join(s for s in ocr_slides if s)
    if len(ocr_text) > len(pptx_text) * 1.2:
        return ocr_text or None
    return pptx_text or ocr_text or None


def _supports_document_text_extraction(filename: str, mime_type: str | None) -> bool:
    ext = os.path.splitext((filename or "").lower())[1]
    normalized_mime = (mime_type or "").lower()
    if ext == ".pdf" or normalized_mime == "application/pdf":
        return True
    if ext in (".docx", ".doc") or normalized_mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return True
    if ext in (".xlsx", ".xls") or normalized_mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return True
    if ext == ".pptx" or normalized_mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return True
    if ext in (".txt", ".md", ".csv", ".log", ".json", ".yaml", ".yml"):
        return True
    return bool(normalized_mime and normalized_mime.startswith("text/"))

